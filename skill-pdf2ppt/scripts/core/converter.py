
import logging
import os
import requests
import tempfile
import time
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import shutil
import uuid
import json
import zipfile
import io
import math
from pathlib import Path

# Import helpers from utils
try:
    from .utils import (
        get_pdf_hash, ensure_cache_dir, recursive_blocks, 
        calculate_font_size_gemini_style, estimate_font_size_by_area, 
        clean_latex_symbols, parse_html_table, is_watermark_element
    )
except ImportError:
    # Fallback for direct execution
    from utils import (
        get_pdf_hash, ensure_cache_dir, recursive_blocks, 
        calculate_font_size_gemini_style, estimate_font_size_by_area, 
        clean_latex_symbols, parse_html_table, is_watermark_element
    )

class MinerUClient:
    def __init__(self, token):
        self.token = token
        self.base_url = "https://mineru.net/api/v4"
        self.headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {self.token}"
        }

    def upload_and_extract(self, file_path, model_version="vlm"):
        """上传 PDF 并等待自动解析"""
        file_name = os.path.basename(file_path)
        data_id = str(uuid.uuid4())
        
        logging.info(f"为文件 {file_name} 申请上传 URL...")
        batch_url = f"{self.base_url}/file-urls/batch"
        data = {
            "files": [{"name": file_name, "data_id": data_id}],
            "model_version": model_version
        }
        
        response = requests.post(batch_url, headers=self.headers, json=data)
        response.raise_for_status()
        result = response.json()
        
        if result.get("code") != 0:
            raise Exception(f"申请上传 URL 失败: {result.get('msg')}")
        
        batch_id = result["data"]["batch_id"]
        upload_url = result["data"]["file_urls"][0]
        
        logging.info(f"✓ 获得 Batch ID: {batch_id}")
        
        # 上传文件
        with open(file_path, 'rb') as f:
            upload_response = requests.put(upload_url, data=f, timeout=120)
            upload_response.raise_for_status()
            
        logging.info(f"✓ 文件上传成功")
        return batch_id

    def get_batch_result(self, batch_id):
        """轮询任务结果并解析 ZIP"""
        result_url = f"{self.base_url}/extract-results/batch/{batch_id}"
        logging.info(f"轮询中...")
        
        poll_count = 0
        max_polls = 180
        
        while poll_count < max_polls:
            response = requests.get(result_url, headers={"Authorization": f"Bearer {self.token}"})
            response.raise_for_status()
            result = response.json()
            
            if result.get("code") != 0:
                raise Exception(f"查询失败: {result.get('msg')}")
            
            extract_results = result.get("data", {}).get("extract_result", [])
            if not extract_results:
                raise Exception("未找到解析结果")
            
            file_result = extract_results[0]
            state = file_result.get("state")
            
            poll_count += 1
            
            if state == "done":
                full_zip_url = file_result.get("full_zip_url")
                if not full_zip_url:
                    raise Exception("未返回结果 URL")
                
                logging.info(f"✓ 解析完成，下载中...")
                zip_response = requests.get(full_zip_url, timeout=120)
                zip_response.raise_for_status()
                
                return zip_response.content
                
            elif state == "failed":
                raise Exception(f"解析失败: {file_result.get('err_msg')}")
            
            if poll_count % 3 == 1:  # 每30秒显示一次
                logging.info(f"  状态: {state} ({poll_count}/{max_polls})")
            time.sleep(10)
        
        raise Exception(f"任务超时")

def save_mineru_result_to_cache(page_num, zip_content, pdf_hash, cache_dir):
    """保存 MinerU 结果到缓存"""
    ensure_cache_dir(cache_dir)
    
    page_dir = os.path.join(cache_dir, f"{pdf_hash}_page_{page_num}")
    if not os.path.exists(page_dir):
        os.makedirs(page_dir)
    
    # 1. 保存 ZIP
    zip_path = os.path.join(page_dir, "mineru_result.zip")
    with open(zip_path, 'wb') as f:
        f.write(zip_content)
    logging.info(f"  ✓ 保存 ZIP: {zip_path}")
    
    # 2. 解析并保存 JSON 和图片
    zip_data = io.BytesIO(zip_content)
    with zipfile.ZipFile(zip_data, 'r') as zip_ref:
        # 保存所有 JSON 文件
        for file_name in zip_ref.namelist():
            if file_name.endswith('.json'):
                json_content = zip_ref.read(file_name)
                json_path = os.path.join(page_dir, os.path.basename(file_name))
                with open(json_path, 'wb') as f:
                    f.write(json_content)
                logging.info(f"  ✓ 保存 JSON: {json_path}")
        
        # 保存图片  
        images_dir = os.path.join(page_dir, "images")
        if not os.path.exists(images_dir):
            os.makedirs(images_dir)
        
        for file_name in zip_ref.namelist():
            if file_name.startswith('images/') and not file_name.endswith('/'):
                img_data = zip_ref.read(file_name)
                img_path = os.path.join(images_dir, os.path.basename(file_name))
                with open(img_path, 'wb') as f:
                    f.write(img_data)
                logging.info(f"  ✓ 保存图片: {img_path}")
    
    logging.info(f"✓ 缓存已保存到: {page_dir}")

def load_mineru_result_from_cache(page_num, pdf_hash, cache_dir):
    """从缓存加载 MinerU 结果"""
    page_dir = os.path.join(cache_dir, f"{pdf_hash}_page_{page_num}")
    zip_path = os.path.join(page_dir, "mineru_result.zip")
    
    if os.path.exists(zip_path):
        logging.info(f"  ✓ 从缓存加载: {page_dir}")
        with open(zip_path, 'rb') as f:
            return f.read()
    
    return None

def parse_mineru_zip(zip_content):
    """解析 MinerU ZIP（使用 content_list.json 和 layout.json）"""
    zip_data = io.BytesIO(zip_content)
    result = {
        'elements': [],
        'images_data': {},
        'page_size': None
    }
    
    with zipfile.ZipFile(zip_data, 'r') as zip_ref:
        # 1. 优先读取 content_list.json
        content_list_files = [f for f in zip_ref.namelist() 
                             if f.endswith('_content_list.json')]
        
        content_list_page_size = None
        if content_list_files:
            json_content = zip_ref.read(content_list_files[0])
            raw_elements = json.loads(json_content)
            # 使用递归展平逻辑，确保存储在 content_list 中的嵌套结构也能被读取
            result['elements'] = recursive_blocks(raw_elements)
            logging.debug(f"读取 {len(result['elements'])} 个元素 (Recursive)")
            
            # 尝试从第一个元素的bbox推断content_list的page_size
            for elem in result['elements']:
                bbox = elem.get('bbox')
                if bbox and len(bbox) == 4:
                    if content_list_page_size is None:
                        content_list_page_size = [0, 0]
                    content_list_page_size[0] = max(content_list_page_size[0], bbox[2])
                    content_list_page_size[1] = max(content_list_page_size[1], bbox[3])
        
        # 2. 读取 layout.json 提取 image_caption 的 bbox
        layout_files = [f for f in zip_ref.namelist() if f.endswith('layout.json')]
        if layout_files:
            layout_content = zip_ref.read(layout_files[0])
            layout_data = json.loads(layout_content)
            
            # 获取layout.json的page_size
            layout_page_size = None
            if 'pdf_info' in layout_data and len(layout_data['pdf_info']) > 0:
                layout_page_size = layout_data['pdf_info'][0].get('page_size')
                para_blocks = layout_data['pdf_info'][0].get('para_blocks', [])
                
                # 计算坐标转换比例
                scale_x = 1.0
                scale_y = 1.0
                if layout_page_size and content_list_page_size:
                    scale_x = content_list_page_size[0] / layout_page_size[0]
                    scale_y = content_list_page_size[1] / layout_page_size[1]
                
                for block in para_blocks:
                    if block.get('type') == 'image':
                        sub_blocks = block.get('blocks', [])
                        for sub_block in sub_blocks:
                            if sub_block.get('type') == 'image_caption':
                                lines = sub_block.get('lines', [])
                                if lines:
                                    for line in lines:
                                        spans = line.get('spans', [])
                                        for span in spans:
                                            if span.get('type') == 'text':
                                                caption_text = span.get('content', '')
                                                caption_bbox = span.get('bbox', [])
                                                
                                                if caption_text and len(caption_bbox) == 4:
                                                    transformed_bbox = [
                                                        int(caption_bbox[0] * scale_x),
                                                        int(caption_bbox[1] * scale_y),
                                                        int(caption_bbox[2] * scale_x),
                                                        int(caption_bbox[3] * scale_y)
                                                    ]
                                                    
                                                    caption_element = {
                                                        'type': 'text',
                                                        'text': caption_text,
                                                        'bbox': transformed_bbox,
                                                        'page_idx': 0,
                                                        '_from_image_caption': True
                                                    }
                                                    result['elements'].append(caption_element)
        
        # 3. 提取图片
        image_files = [f for f in zip_ref.namelist() 
                      if f.startswith('images/') and not f.endswith('/')]
        
        for img_path in image_files:
            img_name = os.path.basename(img_path)
            result['images_data'][img_name] = zip_ref.read(img_path)
    
    return result

def add_image_to_slide(slide, img_path, ppt_x, ppt_y, ppt_width, ppt_height, images_data, temp_dir):
    """辅助函数：添加图片到幻灯片"""
    img_filename = os.path.basename(img_path)
    if img_filename in images_data:
        temp_img_path = os.path.join(temp_dir, img_filename)
        try:
            with open(temp_img_path, 'wb') as f:
                f.write(images_data[img_filename])
            
            slide.shapes.add_picture(
                temp_img_path, 
                ppt_x, 
                ppt_y, 
                width=ppt_width,
                height=ppt_height
            )
            return True
        except Exception as e:
            logging.warning(f"  图片添加失败: {img_filename}, {e}")
            return False
        finally:
            if os.path.exists(temp_img_path):
                try:
                    os.unlink(temp_img_path)
                except:
                    pass
    return False

def create_ppt_table(slide, table_data, ppt_x, ppt_y, ppt_width, ppt_height, page_height, slide_height_inches):
    """
    在PPT中创建原生表格，支持居中对齐和动态字号
    """
    if not table_data or len(table_data) == 0:
        return False
    
    try:
        rows = len(table_data)
        cols = max(len(row) for row in table_data) if table_data else 0
        
        if rows == 0 or cols == 0:
            return False
        
        # 创建表格
        table_shape = slide.shapes.add_table(rows, cols, ppt_x, ppt_y, ppt_width, ppt_height)
        table = table_shape.table
        
        # 计算表格单元格的平均高度
        cell_height_emu = ppt_height / rows
        cell_height_inches = cell_height_emu / 914400
        cell_height_pt = cell_height_inches * 72
        
        # 字号 = 单元格高度的60%
        estimated_font_size = cell_height_pt * 0.6
        font_size = max(10, min(16, int(estimated_font_size)))

        # 填充数据
        for i, row_data in enumerate(table_data):
            for j, cell_data in enumerate(row_data):
                if j < cols:
                    cell = table.cell(i, j)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    cell.text = str(cell_data)
                    
                    text_frame = cell.text_frame
                    text_frame.word_wrap = True
                    
                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
                            run.font.name = 'Microsoft YaHei'
                            if i == 0:
                                run.font.bold = True
        
        logging.debug(f"  表格创建成功: {rows}行×{cols}列, 字号={font_size}pt")
        return True
    except Exception as e:
        logging.warning(f"PPT表格创建失败: {e}")
        return False

def process_elements(elements, slide, images_data, page_width, page_height, 
                       slide_width_emu, slide_height_emu, slide_height_inches, temp_dir,
                       remove_watermark=True):
    """
    处理所有元素并添加到 Slide
    """
    stats = {'文本': 0, '标题': 0, '页脚': 0, '图片': 0, '表格': 0, '列表': 0, '其他': 0, '跳过': 0}
    
    for i, element in enumerate(elements):
        if not isinstance(element, dict):
            stats['跳过'] += 1
            continue
        
        # 去水印检测
        if remove_watermark and is_watermark_element(element, elements, page_width, page_height):
            stats['跳过'] += 1
            continue
        
        elem_type = element.get('type')
        bbox = element.get('bbox')
        
        if not bbox or len(bbox) != 4:
            logging.debug(f"  元素 {i}: type={elem_type}, bbox=无效")
            stats['跳过'] += 1
            continue
        
        # 这里假设 content_list.json 的 bbox 是像素坐标
        x1_px, y1_px, x2_px, y2_px = bbox
        
        # 转换为 PPT EMU 坐标
        x1_ratio = x1_px / page_width
        y1_ratio = y1_px / page_height
        x2_ratio = x2_px / page_width
        y2_ratio = y2_px / page_height
        
        ppt_x = Emu(int(x1_ratio * slide_width_emu))
        ppt_y = Emu(int(y1_ratio * slide_height_emu))
        ppt_width = Emu(int((x2_ratio - x1_ratio) * slide_width_emu))
        ppt_height = Emu(int((y2_ratio - y1_ratio) * slide_height_emu))
        
        # 最小尺寸保护
        min_width = Emu(int(0.01 * slide_width_emu))
        min_height = Emu(int(0.01 * slide_height_emu))
        ppt_width = max(ppt_width, min_width)
        ppt_height = max(ppt_height, min_height)
        
        # 1. 处理图片
        if elem_type == 'image':
            img_path = element.get('img_path')
            if img_path:
                if add_image_to_slide(slide, img_path, ppt_x, ppt_y, ppt_width, ppt_height, images_data, temp_dir):
                    stats['图片'] += 1
                else:
                    stats['跳过'] += 1
            else:
                stats['跳过'] += 1

        # 2. 处理表格
        elif elem_type == 'table':
            table_created = False
            
            # 策略1: 优先尝试HTML表格
            table_body = element.get('table_body')
            if table_body:
                table_data = parse_html_table(table_body)
                if table_data:
                    if create_ppt_table(slide, table_data, ppt_x, ppt_y, ppt_width, ppt_height, page_height, slide_height_inches):
                        stats['表格'] += 1
                        table_created = True
                        
                        # 表格标题
                        table_caption = element.get('table_caption')
                        if table_caption and isinstance(table_caption, list) and len(table_caption) > 0:
                            caption_text = '\n'.join(table_caption)
                            caption_text = clean_latex_symbols(caption_text)
                            caption_y = ppt_y + ppt_height
                            caption_height = Emu(int(0.05 * slide_height_emu))
                            try:
                                caption_box = slide.shapes.add_textbox(ppt_x, caption_y, ppt_width, caption_height)
                                caption_p = caption_box.text_frame.paragraphs[0]
                                caption_p.text = caption_text
                                caption_p.font.size = Pt(10)
                                caption_p.font.name = 'Microsoft YaHei'
                                caption_p.font.italic = True
                                stats['文本'] += 1
                            except:
                                pass
            
            # 策略2: 降级到图片
            if not table_created:
                img_path = element.get('img_path')
                if img_path and add_image_to_slide(slide, img_path, ppt_x, ppt_y, ppt_width, ppt_height, images_data, temp_dir):
                    stats['表格'] += 1
                    table_created = True
                    
                    # 表格标题
                    table_caption = element.get('table_caption')
                    if table_caption and isinstance(table_caption, list) and len(table_caption) > 0:
                        caption_text = '\n'.join(table_caption)
                        caption_text = clean_latex_symbols(caption_text)
                        caption_y = ppt_y + ppt_height
                        caption_height = Emu(int(0.05 * slide_height_emu))
                        try:
                            caption_box = slide.shapes.add_textbox(ppt_x, caption_y, ppt_width, caption_height)
                            caption_p = caption_box.text_frame.paragraphs[0]
                            caption_p.text = caption_text
                            caption_p.font.size = Pt(10)
                            caption_p.font.name = 'Microsoft YaHei'
                            caption_p.font.italic = True
                            stats['文本'] += 1
                        except:
                            pass

            # 策略3: 降级到纯文本
            if not table_created:
                table_text = element.get('text') or element.get('content', '')
                if not table_text: table_text = "[表格数据]"
                
                try:
                    txBox = slide.shapes.add_textbox(ppt_x, ppt_y, ppt_width, ppt_height)
                    p = txBox.text_frame.paragraphs[0]
                    p.text = table_text
                    p.font.size = Pt(10)
                    p.font.name = 'Consolas'
                    stats['表格'] += 1
                except:
                    stats['跳过'] += 1
        
        # 3. 处理文字 & 列表
        elif elem_type in ['text', 'title', 'footer', 'list']:
            text_content = ""
            is_title = (elem_type == 'title')
            
            if elem_type == 'list':
                list_items = element.get('list_items') or []
                if list_items:
                    cleaned_items = [clean_latex_symbols(item) for item in list_items]
                    text_content = '\n'.join(cleaned_items)
            else:
                text_content = element.get('text') or element.get('content', '')
                text_content = clean_latex_symbols(text_content)

            if not text_content:
                stats['跳过'] += 1
                continue
            
            try:
                txBox = slide.shapes.add_textbox(ppt_x, ppt_y, ppt_width, ppt_height)
                tf = txBox.text_frame
                
                char_count = len(text_content.strip())
                if char_count <= 6:
                    tf.word_wrap = False
                else:
                    tf.word_wrap = True
                
                p = tf.paragraphs[0]
                p.text = text_content
                
                bbox_w = x2_px - x1_px
                bbox_h = y2_px - y1_px
                
                estimated_fontSize_px = estimate_font_size_by_area(bbox_w, bbox_h, char_count, is_title)
                font_size = calculate_font_size_gemini_style(estimated_fontSize_px, page_height, slide_height_inches)
                
                if elem_type == 'title':
                    font_size = max(font_size, 20)
                    p.font.bold = True
                elif elem_type == 'footer':
                    font_size = min(font_size, 10)
                elif elem_type == 'list':
                    font_size = min(font_size, 20)
                
                p.font.size = Pt(font_size)
                p.font.name = 'Microsoft YaHei'
                
                stat_key_map = {'text': '文本', 'title': '标题', 'footer': '页脚', 'list': '列表'}
                stats[stat_key_map.get(elem_type, elem_type)] += 1
            except Exception as e:
                logging.warning(f"  文本添加失败: {e}")
                stats['跳过'] += 1
        else:
            stats['其他'] += 1

    return stats

def convert_pdf_to_ppt(pdf_input_path, ppt_output_path, mineru_token, 
                         ppt_slide_width=16, ppt_slide_height=9, 
                         use_cache=True, cache_dir="temp",
                         remove_watermark=True):
    """
    核心转换函数
    """
    if not os.path.exists(pdf_input_path):
        logging.error(f"输入 PDF 不存在: {pdf_input_path}")
        raise FileNotFoundError(f"文件不存在: {pdf_input_path}")
    
    if not mineru_token:
        logging.error("MINERU_TOKEN 未设置")
        raise ValueError("Token不能为空")

    pdf_hash = get_pdf_hash(pdf_input_path)
    logging.info(f"PDF 哈希: {pdf_hash}")

    temp_dir = tempfile.mkdtemp()
    
    try:
        client = MinerUClient(mineru_token)
        tasks = []
        
        logging.info(f"正在拆分 PDF '{pdf_input_path}'...")
        doc = fitz.open(pdf_input_path)
        
        logging.info(f"正在批量上传 {len(doc)} 个页面以并发处理...")
        
        for i, page in enumerate(doc):
            page_num = i + 1
            
            cached_zip = None
            if use_cache:
                cached_zip = load_mineru_result_from_cache(page_num, pdf_hash, cache_dir)
            
            if cached_zip:
                logging.info(f"  [页 {page_num}] 命中本地缓存")
                tasks.append({
                    'type': 'cached',
                    'page_num': page_num,
                    'content': cached_zip
                })
            else:
                single_page_doc = fitz.open()
                single_page_doc.insert_pdf(doc, from_page=i, to_page=i)
                page_path = os.path.join(temp_dir, f"page_{page_num}.pdf")
                single_page_doc.save(page_path)
                
                try:
                    batch_id = client.upload_and_extract(page_path)
                    tasks.append({
                        'type': 'api',
                        'page_num': page_num,
                        'batch_id': batch_id,
                        'page_path': page_path
                    })
                except Exception as e:
                    logging.error(f"  [页 {page_num}] 上传失败: {e}")
                    tasks.append({
                        'type': 'error',
                        'page_num': page_num,
                        'error': str(e)
                    })

        doc.close()

        all_pages_results = []
        logging.info(f"\n{'-'*40}\n所有页面已提交，开始轮询结果\n{'-'*40}")
        
        for task in tasks:
            page_num = task['page_num']
            
            if task['type'] == 'error':
                logging.warning(f"[页 {page_num}] 跳过 (上传阶段失败)")
                all_pages_results.append(None)
                continue
                
            if task['type'] == 'cached':
                zip_content = task['content']
            else:
                batch_id = task['batch_id']
                page_path = task.get('page_path')
                zip_content = None
                
                max_retries = 3
                for attempt in range(max_retries + 1):
                    try:
                        if attempt > 0:
                            logging.info(f"[页 {page_num}] 第 {attempt} 次重试: 重新上传并解析...")
                            if page_path and os.path.exists(page_path):
                                batch_id = client.upload_and_extract(page_path)
                            else:
                                raise Exception("无法重试: 临时文件丢失")
                        
                        logging.info(f"[页 {page_num}] 正在等待任务 {batch_id} ...")
                        zip_content = client.get_batch_result(batch_id)
                        
                        if use_cache:
                            save_mineru_result_to_cache(page_num, zip_content, pdf_hash, cache_dir)
                        break
                        
                    except Exception as e:
                        if attempt < max_retries:
                            logging.warning(f"[页 {page_num}] 解析失败: {e}。等待 5 秒后重试...")
                            time.sleep(5)
                        else:
                            logging.error(f"[页 {page_num}] 重试 {max_retries} 次后最终失败: {e}")
                            zip_content = None
                
                if zip_content is None:
                    all_pages_results.append(None)
                    continue

            page_result = parse_mineru_zip(zip_content)
            
            if page_result and page_result.get('elements'):
                max_x = max_y = 0
                for elem in page_result['elements']:
                    if elem.get('bbox') and len(elem['bbox']) == 4:
                        max_x = max(max_x, elem['bbox'][2])
                        max_y = max(max_y, elem['bbox'][3])
                
                page_width = max_x if max_x > 0 else 1000
                page_height = max_y if max_y > 0 else 1000
                
                all_pages_results.append({
                    'elements': page_result['elements'],
                    'images_data': page_result.get('images_data', {}),
                    'page_width': page_width,
                    'page_height': page_height
                })
            else:
                logging.warning(f"[页 {page_num}] 未返回有效内容")
                all_pages_results.append(None)

        logging.info(f"\n{'='*40}\n开始合成 PPT\n{'='*40}")
        prs = Presentation()
        prs.slide_width = Inches(ppt_slide_width)
        prs.slide_height = Inches(ppt_slide_height)
        slide_width_emu = prs.slide_width
        slide_height_emu = prs.slide_height

        for i, page_data in enumerate(all_pages_results):
            slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(slide_layout)
            
            if page_data:
                stats = process_elements(
                    page_data['elements'],
                    slide,
                    page_data['images_data'],
                    page_data['page_width'],
                    page_data['page_height'],
                    slide_width_emu,
                    slide_height_emu,
                    ppt_slide_height,
                    temp_dir,
                    remove_watermark
                )
                stats_str = ', '.join([f"{k}{v}" for k, v in stats.items() if v > 0])
                logging.info(f"生成幻灯片 {i+1}（{stats_str}）")
            else:
                txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
                txBox.text_frame.text = f"Page {i+1} 解析失败或无内容"

        prs.save(ppt_output_path)
        logging.info(f"\n✅ 完成: {ppt_output_path}")
        return ppt_output_path

    except Exception as e:
        logging.error(f"失败: {e}", exc_info=True)
        raise e
    finally:
        if os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
